import os

# Force fully offline / no-telemetry operation. Must be set before
# huggingface_hub, sentence_transformers, or chromadb are imported below -
# they read these at import time to decide whether to reach out to the
# network. This is what keeps embedded document content from ever leaving
# the machine: no model-update checks, no anonymized usage pings.
os.environ.setdefault("HF_HUB_OFFLINE", "1")
os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")
os.environ.setdefault("ANONYMIZED_TELEMETRY", "False")

import json
import time
import uuid
import logging
import requests
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Any, AsyncGenerator, Optional

import httpx
from chromadb.config import Settings as ChromaSettings
from PIL import Image
import pytesseract
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from dotenv import load_dotenv
from sympy import re

load_dotenv(Path(__file__).parent / ".env")

# ── Timing instrumentation ──────────────────────────────────────────────────
# Dedicated logger so pipeline timing can be filtered from app logs
# (e.g. `grep TIMING`) without touching the root logger's config.
timing_logger = logging.getLogger("rag_timing")
timing_logger.setLevel(logging.INFO)
if not timing_logger.handlers:
    _handler = logging.StreamHandler()
    _handler.setFormatter(logging.Formatter(
        "%(asctime)s.%(msecs)03d [TIMING] %(message)s", datefmt="%H:%M:%S"
    ))
    timing_logger.addHandler(_handler)
    timing_logger.propagate = False


class StageTimer:
    """Accumulates named stage durations for one request and logs a summary."""

    def __init__(self, request_id: str, label: str):
        self.request_id = request_id
        self.label = label
        self.stages: List[tuple] = []
        self._t_start = time.perf_counter()

    def stage(self, name: str):
        return _StageContext(self, name)

    def _record(self, name: str, ms: float):
        self.stages.append((name, ms))
        timing_logger.info(f"[{self.request_id}] {self.label}.{name}: {ms:.1f}ms")

    def summary(self):
        total_ms = (time.perf_counter() - self._t_start) * 1000
        breakdown = ", ".join(f"{n}={ms:.1f}ms" for n, ms in self.stages)
        timing_logger.info(
            f"[{self.request_id}] {self.label}.TOTAL: {total_ms:.1f}ms ({breakdown})"
        )
        return total_ms


class _StageContext:
    def __init__(self, timer: StageTimer, name: str):
        self.timer = timer
        self.name = name

    def __enter__(self):
        self._t0 = time.perf_counter()
        return self

    def __exit__(self, exc_type, exc, tb):
        ms = (time.perf_counter() - self._t0) * 1000
        self.timer._record(self.name, ms)
        return False

OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "phi3.5:3.8b-mini-instruct-q4_K_M")
# This is a dedicated app server, not a personal laptop that needs its RAM
# back - keep the model resident all day instead of unloading after 30m of
# inactivity, so only the very first message of the day pays the cold-load
# cost (and _warm_model_async even avoids that one).
OLLAMA_KEEP_ALIVE = os.getenv("OLLAMA_KEEP_ALIVE", "24h")
CHROMA_DIR = str(Path(__file__).parent / "chroma_db")

GENERATION_OPTIONS = {
    # Ceiling only - raising this doesn't cost latency by itself, it just
    # avoids silent truncation as the corpus/history grows (was 4096, too
    # tight once the collection grew past a couple of short documents).
    # What actually drives latency is how many tokens get put IN the prompt
    # (retrieval k, history length - see retrieve() and chat_stream), not
    # this number.
    "num_ctx": 8192,
    # qwen2.5-instruct (unlike deepseek-r1) doesn't spend budget on hidden
    # <think> reasoning, so a much smaller cap than before is safe. Bounds
    # worst-case generation time without truncating typical answers. Raised
    # from 250 - the prompt now asks for bullet lists, bold terms, and
    # citations, which use more tokens per answer than plain prose did.
    "num_predict": 400,
    "temperature": 0.2,
    "top_p": 0.5,
    "top_k": 20,
    "repeat_penalty": 1.2,
}

SUPPORTED_LOADERS = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
    ".txt": TextLoader,
    ".md": TextLoader,
    ".csv": TextLoader,
    ".json": TextLoader,
    ".rtf": TextLoader,
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
OCR_EXTENSIONS = IMAGE_EXTENSIONS | {".pptx"}


def _ocr_image(image: Image.Image) -> str:
    return pytesseract.image_to_string(image).strip()


def _load_image(file_path: str) -> List[Document]:
    text = _ocr_image(Image.open(file_path))
    return [Document(page_content=text)]


def _load_pptx(file_path: str) -> List[Document]:
    prs = Presentation(file_path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_text.append(shape.text_frame.text.strip())
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = Image.open(BytesIO(shape.image.blob))
                    ocr_text = _ocr_image(image)
                    if ocr_text:
                        slide_text.append(ocr_text)
                except Exception:
                    pass
        if slide_text:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return [Document(page_content="\n\n".join(parts))]


EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"


class LocalEmbeddings(Embeddings):
    """CPU-resident embedding model - loaded once, never evicted.

    Previously embeddings went through Ollama (nomic-embed-text). On memory-
    constrained hosts, Ollama can only keep one model loaded at a time, so
    every retrieval call evicted the LLM to load the embedder and vice versa,
    adding tens of seconds of reload latency to every single request. Running
    the embedding model in-process removes it from Ollama's rotation entirely,
    so Ollama only ever has to manage the one generation model.
    """

    def __init__(self, model_name: str = EMBEDDING_MODEL_NAME):
        from sentence_transformers import SentenceTransformer
        self._model = SentenceTransformer(model_name, device="cpu")

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        return self._model.encode(texts, convert_to_numpy=True, show_progress_bar=False).tolist()

    def embed_query(self, text: str) -> List[float]:
        return self._model.encode(text, convert_to_numpy=True, show_progress_bar=False).tolist()


class RAGEngine:
    def __init__(self):
        os.makedirs(CHROMA_DIR, exist_ok=True)

        self.embeddings = LocalEmbeddings()
        # The embedding model's first real .encode() call pays a one-time
        # backend/JIT warmup cost independent of Ollama's own warmup
        # (measured 20s+ on first use) - eating that here, synchronously,
        # means it happens once at startup instead of on someone's first
        # question. Cheap enough (short string) not to delay startup much.
        self.embeddings.embed_query("warmup")
        print(f"✓ Embedding model ({EMBEDDING_MODEL_NAME}, local) ready")

        self.vectorstore = Chroma(
            persist_directory=CHROMA_DIR,
            embedding_function=self.embeddings,
            collection_metadata={"hnsw:space": "cosine"},
            client_settings=ChromaSettings(anonymized_telemetry=False),
        )
        print(f"✓ RAG engine ready (model: {OLLAMA_MODEL})")
        self._warm_model_async()

    def _warm_model_async(self):
        """Load the model into Ollama's memory in the background at startup.

        Without this, whoever sends the first chat message pays the full
        disk-load cost (can be 60s+) and the app looks broken/unresponsive
        - Ollama only loads a model into RAM on its first request, not when
        it starts serving. Fire-and-forget: failures here just mean the
        first real message pays the cold-load cost instead, same as before.
        """
        import threading

        def _warm():
            try:
                requests.post(
                    f"{OLLAMA_URL}/api/generate",
                    json={"model": OLLAMA_MODEL, "prompt": "hi", "stream": False, "keep_alive": OLLAMA_KEEP_ALIVE},
                    timeout=180,
                )
                print(f"✓ Model warmed up ({OLLAMA_MODEL})")
            except Exception as e:
                print(f"Warning: model warmup failed, first real message will be slow instead ({e})")

        threading.Thread(target=_warm, daemon=True).start()

    # ── Retrieval ─────────────────────────────────────────────────────────────

    def retrieve(
        self, question: str, department: Optional[str] = None, request_id: str = "-"
    ) -> List[Document]:
        # Each retrieved chunk is now ~800-1000 coherent chars (vs ~350
        # before), so fewer chunks carry as much real content while costing
        # far less prompt-processing time - a direct latency lever.
        doc_count = 3 if len(question.split()) < 10 else 4

        # Embedding generation and vector search are timed separately (rather
        # than via retriever.invoke(), which bundles both) so it's possible to
        # tell whether latency comes from the embedding model or the ANN search.
        t0 = time.perf_counter()
        query_embedding = self.embeddings.embed_query(question)
        t1 = time.perf_counter()
        timing_logger.info(
            f"[{request_id}] retrieve.embedding_generation: {(t1 - t0) * 1000:.1f}ms "
            f"(model={EMBEDDING_MODEL_NAME}, question_len={len(question)})"
        )

        # department=None means "search everything" (admin/manager "all"
        # mode). Otherwise chunks are hard-filtered to that department at the
        # vector-search level - this is what actually keeps one department's
        # documents from ever surfacing in another department's answers,
        # not just a UI-level hide.
        search_filter = {"department": department} if department else None

        # fetch_k widened (5x vs previous 3x) - with multiple documents of very
        # different sizes sharing one collection (a 220-page novel dwarfs the
        # 2 short business docs in chunk count), MMR needs a bigger candidate
        # pool or it barely ever surfaces the small documents.
        docs = self.vectorstore.max_marginal_relevance_search_by_vector(
            embedding=query_embedding, k=doc_count, fetch_k=doc_count * 5, filter=search_filter,
        )
        t2 = time.perf_counter()
        timing_logger.info(
            f"[{request_id}] retrieve.vector_search_mmr: {(t2 - t1) * 1000:.1f}ms "
            f"(k={doc_count}, fetch_k={doc_count * 3}, department={department}, docs_returned={len(docs)})"
        )

        return docs


    def _clean_response(self, text: str) -> str:
        """Remove filler and make response tighter."""
        # Remove "sure", "certainly", "of course" openers
        text = re.sub(r'^(sure|certainly|of course|based on|according to)[\s,]*', '', text, flags=re.IGNORECASE)
        # Remove redundant quotes
        text = re.sub(r'"\s*(.+?)\s*"', r'\1', text)
        # Collapse multiple spaces
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    # ── Synchronous generation (for /api/chat) ─────────────────────────────

    def _generate_sync(self, prompt: str, request_id: str = "-") -> str:
        t0 = time.perf_counter()
        try:
            resp = requests.post(
                f"{OLLAMA_URL}/api/generate",
                json={
                    "model": OLLAMA_MODEL, "prompt": prompt, "stream": False,
                    "options": GENERATION_OPTIONS, "keep_alive": OLLAMA_KEEP_ALIVE,
                },
                timeout=300,
            )
            resp.raise_for_status()
            result = resp.json().get("response", "").strip()
            elapsed_ms = (time.perf_counter() - t0) * 1000
            timing_logger.info(
                f"[{request_id}] generate.model_inference: {elapsed_ms:.1f}ms "
                f"(model={OLLAMA_MODEL}, prompt_chars={len(prompt)}, "
                f"num_predict={GENERATION_OPTIONS['num_predict']}, num_ctx={GENERATION_OPTIONS['num_ctx']})"
            )
            return result
        except Exception as e:
            elapsed_ms = (time.perf_counter() - t0) * 1000
            timing_logger.info(f"[{request_id}] generate.model_inference: FAILED after {elapsed_ms:.1f}ms ({e})")
            return f"[Generation error: {e}]"

    def query_text(self, question: str, department: Optional[str] = None) -> Dict[str, Any]:
        request_id = uuid.uuid4().hex[:8]
        timer = StageTimer(request_id, "query_text")
        timing_logger.info(f"[{request_id}] query_text.start: question_chars={len(question)}")

        with timer.stage("retrieval"):
            docs = self.retrieve(question, department=department, request_id=request_id)

        if not docs:
            timer.summary()
            return {"answer": "No relevant information found in the documents.", "sources": [], "confidence": 0.0}

        with timer.stage("document_processing"):
            context = "\n".join(
                f"[{i+1}] {d.page_content[:250]}"  # Truncate each source to 250 chars
                for i, d in enumerate(docs)
            )
            sources = [
                {
                    "text": d.page_content[:300],
                    "source": Path(d.metadata.get("source", "unknown")).name,
                    "page": d.metadata.get("page", 0) + 1 if d.metadata.get("page") is not None else 1,
                }
                for d in docs
            ]

        with timer.stage("prompt_construction"):
            prompt = f"""You are a precise knowledge assistant. Answer ONLY using facts from the context below.

RULES:
- Answer directly without preamble
- If the context doesn't contain the answer, say "I don't have this information"
- Cite the source number [1] [2] if referencing specific documents
- Keep answers concise and factual
- Never make assumptions or use outside knowledge

Context:
{context}

Question: {question}
Answer:"""

        with timer.stage("model_inference"):
            answer = self._generate_sync(prompt, request_id=request_id)

        with timer.stage("response_cleanup"):
            answer = self._clean_response(answer)

        confidence = 0.9 if len(answer) > 20 and "not found" not in answer.lower() else 0.3
        timer.summary()
        return {"answer": answer, "sources": sources, "confidence": confidence}

    # ── Async streaming generation (for /api/chat/stream) ─────────────────

    async def generate_stream(self, prompt: str, request_id: str = "-") -> AsyncGenerator[str, None]:
        t_request_sent = time.perf_counter()
        first_token_at = None
        token_count = 0

        async with httpx.AsyncClient(timeout=300.0) as client:
            async with client.stream(
                "POST",
                f"{OLLAMA_URL}/api/generate",
                json={
                    "model": OLLAMA_MODEL, "prompt": prompt, "stream": True,
                    "options": GENERATION_OPTIONS, "keep_alive": OLLAMA_KEEP_ALIVE,
                },
            ) as resp:
                t_headers_received = time.perf_counter()
                timing_logger.info(
                    f"[{request_id}] generate_stream.connect_and_headers: "
                    f"{(t_headers_received - t_request_sent) * 1000:.1f}ms "
                    f"(model={OLLAMA_MODEL}, prompt_chars={len(prompt)})"
                )
                resp.raise_for_status()
                async for line in resp.aiter_lines():
                    if line:
                        try:
                            data = json.loads(line)
                            if not data.get("done"):
                                token = data.get("response", "")
                                if token:
                                    if first_token_at is None:
                                        first_token_at = time.perf_counter()
                                        timing_logger.info(
                                            f"[{request_id}] generate_stream.time_to_first_token: "
                                            f"{(first_token_at - t_request_sent) * 1000:.1f}ms"
                                        )
                                    token_count += 1
                                    yield token
                        except json.JSONDecodeError:
                            pass

        t_end = time.perf_counter()
        total_ms = (t_end - t_request_sent) * 1000
        generation_ms = (t_end - first_token_at) * 1000 if first_token_at else 0.0
        timing_logger.info(
            f"[{request_id}] generate_stream.model_inference_total: {total_ms:.1f}ms "
            f"(time_to_first_token={((first_token_at - t_request_sent) * 1000) if first_token_at else 0:.1f}ms, "
            f"token_generation={generation_ms:.1f}ms, tokens={token_count}, "
            f"avg_ms_per_token={(generation_ms / token_count) if token_count else 0:.1f})"
        )

    # ── Document indexing ─────────────────────────────────────────────────

    def add_document(self, file_path: str, department: str = "general") -> int:
        path = Path(file_path)
        ext = path.suffix.lower()
        filename = path.name

        if ext not in SUPPORTED_LOADERS and ext not in OCR_EXTENSIONS:
            return 0

        # Deduplication: skip files already in the vector store - scoped to
        # (filename, department) so two departments can each have their own
        # "handbook.pdf" without colliding.
        try:
            existing = self.vectorstore.get(
                where={"$and": [{"filename": filename}, {"department": department}]}
            )
            if existing.get("ids"):
                print(f"Already indexed: {filename} in {department} ({len(existing['ids'])} chunks)")
                return len(existing["ids"])
        except Exception:
            pass

        try:
            if ext in IMAGE_EXTENSIONS:
                docs = _load_image(file_path)
            elif ext == ".pptx":
                docs = _load_pptx(file_path)
            else:
                loader_cls = SUPPORTED_LOADERS[ext]
                loader = loader_cls(file_path) if ext not in (".txt", ".md", ".csv", ".json", ".rtf") else loader_cls(file_path, encoding="utf-8")
                docs = loader.load()

            if not docs or not any(d.page_content.strip() for d in docs):
                raise Exception("Document loaded but contains no content")

            combined = "\n\n".join(d.page_content for d in docs)
            combined_doc = Document(
                page_content=combined,
                metadata={"filename": filename, "file_type": ext, "source": file_path, "department": department},
            )

            splitter = RecursiveCharacterTextSplitter(
                # 400/50 fragmented narrative prose into sub-sentence pieces
                # (a 222-page novel produced ~950 chunks), which made the LLM
                # paraphrase disconnected fragments into garbled answers.
                # 1000/150 keeps whole paragraphs together for coherent context.
                chunk_size=1000,
                chunk_overlap=150,
                separators=["\n\n", "\n", ". ", " ", ""],
            )
            chunks = splitter.split_documents([combined_doc])

            if not chunks:
                raise Exception("No chunks created from document")

            self.vectorstore.add_documents(chunks)
            print(f"✓ Indexed {len(chunks)} chunks to {filename}")
            return len(chunks)

        except Exception as e:
            raise Exception(f"Failed to process {filename}: {e}")

    def delete_document(self, filename: str, department: Optional[str] = None) -> bool:
        try:
            where = (
                {"$and": [{"filename": filename}, {"department": department}]}
                if department else {"filename": filename}
            )
            results = self.vectorstore.get(where=where)
            if not results.get("ids"):
                return False
            self.vectorstore.delete(ids=results["ids"])
            print(f"✓ Deleted {len(results['ids'])} chunks for: {filename}")
            return True
        except Exception as e:
            raise Exception(f"Error deleting {filename}: {e}")

    def sweep_orphaned_chunks(self) -> int:
        """Remove any chunk whose source file no longer exists on disk.

        Deleting a document is normally just delete_document() removing its
        chunks by filename+department, but that's a best-effort match on
        metadata - if it's ever silently skipped (an exception swallowed
        upstream, a crash mid-delete, someone removing a file by hand
        outside the app), the chunks are permanently orphaned in Chroma
        with nothing to point at them: no document to delete again, no UI
        surface showing they exist, just quietly growing storage forever.
        Run at startup as a safety net so drift never survives a restart.
        """
        try:
            results = self.vectorstore.get(include=["metadatas"])
        except Exception as e:
            print(f"Warning: could not scan for orphaned chunks: {e}")
            return 0

        orphan_ids = []
        orphan_files = set()
        for chunk_id, meta in zip(results.get("ids", []), results.get("metadatas", [])):
            source = (meta or {}).get("source")
            if source and not Path(source).exists():
                orphan_ids.append(chunk_id)
                orphan_files.add(meta.get("filename", source))

        if orphan_ids:
            self.vectorstore.delete(ids=orphan_ids)
            print(f"✓ Swept {len(orphan_ids)} orphaned chunks for {len(orphan_files)} missing file(s): {', '.join(sorted(orphan_files))}")
        return len(orphan_ids)

    def delete_department_documents(self, department: str) -> int:
        """Remove every chunk tagged with this department - used when a
        department itself is deleted (cascade)."""
        try:
            results = self.vectorstore.get(where={"department": department})
            if not results.get("ids"):
                return 0
            self.vectorstore.delete(ids=results["ids"])
            print(f"✓ Deleted {len(results['ids'])} chunks for department: {department}")
            return len(results["ids"])
        except Exception as e:
            raise Exception(f"Error deleting department {department}: {e}")